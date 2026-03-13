# Presentation Export Plugin for GYN Oncology Tumor Board
#
# Generates a 3-slide PPTX summary from all agent outputs:
#   Slide 1: Patient Overview (demographics, FIGO, molecular profile)
#   Slide 2: Clinical Findings (pathology, radiology, tumor marker chart)
#   Slide 3: Treatment & Trials (NCCN recs, clinical trials, consensus)
#
# Uses python-pptx with a pre-generated template (tumor_board_slides.pptx).
# Content is summarized via Azure ChatCompletion with SlideContent structured output.

from __future__ import annotations

import json
import logging
import os
from io import BytesIO
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from semantic_kernel.connectors.ai.open_ai.prompt_execution_settings.azure_chat_prompt_execution_settings import (
    AzureChatPromptExecutionSettings,
)
from semantic_kernel.contents.chat_history import ChatHistory
from semantic_kernel.functions import kernel_function

from data_models.chat_artifact import ChatArtifact, ChatArtifactIdentifier
from data_models.chat_context import ChatContext
from data_models.data_access import DataAccess
from data_models.plugin_configuration import PluginConfiguration
from data_models.tumor_board_summary import SlideContent
from routes.patient_data.patient_data_routes import get_chat_artifacts_url
from utils.model_utils import model_supports_temperature

logger = logging.getLogger(__name__)

OUTPUT_PPTX_FILENAME = "tumor_board_slides-{}.pptx"
TEMPLATE_PPTX_FILENAME = "tumor_board_slides.pptx"

# Colors matching the template
NAVY = RGBColor(0x1B, 0x36, 0x5D)
TEAL = RGBColor(0x00, 0x7C, 0x91)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
BULLET_COLOR = RGBColor(0x44, 0x44, 0x44)

SLIDE_SUMMARIZATION_PROMPT = """\
You are a medical presentation writer for a GYN Oncology Tumor Board.
Given the full agent outputs below, create concise slide content.

Rules:
- Each bullet must be ≤20 words, clear, and clinically precise.
- overview_bullets: max 6 items covering demographics, history, FIGO stage, molecular profile, ECOG, key diagnoses.
- findings_bullets: max 6 items covering pathology (histology, grade, IHC, molecular classification), radiology (key imaging findings), and tumor markers.
- treatment_bullets: max 6 items covering NCCN recommendations, board consensus, and follow-up plan.
- trial_entries: max 3 items, format each as "NCT# — Brief title (Phase X)".
- overview_title: "Patient {patient_id} — {cancer_type}"
- overview_subtitle: "FIGO {figo_stage} | {molecular_profile} | {date}"
- If the patient is an outside transfer, include referral source and reason in overview_bullets.
- findings_title: "Pathology & Imaging Findings"
- findings_chart_title: name of primary tumor marker (e.g., "CA-125 Trend")
- treatment_title: "Treatment Plan & Clinical Trials"

Respond with valid JSON matching the SlideContent schema exactly.
"""


def create_plugin(plugin_config: PluginConfiguration):
    return PresentationExportPlugin(
        kernel=plugin_config.kernel,
        chat_ctx=plugin_config.chat_ctx,
        data_access=plugin_config.data_access,
    )


class PresentationExportPlugin:
    def __init__(self, kernel, chat_ctx: ChatContext, data_access: DataAccess):
        # presentation_export.py is at tools/presentation_export.py (2 levels below scenarios/default/)
        self.root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self.kernel = kernel
        self.chat_ctx = chat_ctx
        self.data_access = data_access

    @kernel_function(
        description="Generate a 3-slide PowerPoint presentation summarizing the GYN tumor board case. "
        "Produces slides for Patient Overview, Clinical Findings (with tumor marker chart), "
        "and Treatment Plan with Clinical Trials."
    )
    async def export_to_pptx(
        self,
        patient_age: str,
        patient_gender: str,
        cancer_type: str,
        pathology_findings: str,
        radiology_findings: str,
        treatment_plan: str,
        clinical_trials: str,
        figo_stage: str = "",
        molecular_profile: str = "",
        tumor_markers: str = "",
        surgical_findings: str = "",
        board_discussion: str = "",
        oncologic_history: str = "",
    ) -> str:
        """Generate a 3-slide PPTX tumor board summary.

        Args:
            patient_age: Patient age.
            patient_gender: Patient gender.
            cancer_type: Cancer type (e.g., "High-grade serous ovarian carcinoma").
            pathology_findings: Pathology findings including IHC and molecular.
            radiology_findings: CT/MRI/PET findings summary.
            treatment_plan: NCCN-based treatment recommendation.
            clinical_trials: Eligible clinical trials summary.
            figo_stage: FIGO stage (e.g., "IIIC").
            molecular_profile: Molecular profile (BRCA, HRD, MMR, etc.).
            tumor_markers: Tumor marker trends (CA-125, HE4, etc.) as JSON or text.
            surgical_findings: Surgical/debulking findings.
            board_discussion: Tumor board consensus discussion points.
            oncologic_history: Structured prior oncologic history (diagnosis, treatments, referral reason).

        Returns:
            str: HTML link to download the generated PPTX file.
        """
        patient_id = self.chat_ctx.patient_id
        conversation_id = self.chat_ctx.conversation_id

        # 1. Summarize all agent data into slide-friendly content via LLM
        all_data = {
            "patient_id": patient_id,
            "patient_age": patient_age,
            "patient_gender": patient_gender,
            "cancer_type": cancer_type,
            "figo_stage": figo_stage,
            "molecular_profile": molecular_profile,
            "pathology_findings": pathology_findings,
            "radiology_findings": radiology_findings,
            "tumor_markers": tumor_markers,
            "surgical_findings": surgical_findings,
            "treatment_plan": treatment_plan,
            "clinical_trials": clinical_trials,
            "board_discussion": board_discussion,
            "oncologic_history": oncologic_history,
        }
        slide_content = await self._summarize_for_slides(all_data)

        # 2. Generate tumor marker chart
        marker_chart = self._create_marker_chart(tumor_markers)

        # 3. Load template and build slides
        template_path = os.path.join(self.root_dir, "templates", TEMPLATE_PPTX_FILENAME)
        prs = Presentation(template_path)

        self._build_slide_1(prs.slides[0], slide_content)
        self._build_slide_2(prs.slides[1], slide_content, marker_chart)
        self._build_slide_3(prs.slides[2], slide_content)

        # 4. Save to blob storage
        stream = BytesIO()
        prs.save(stream)

        artifact_id = ChatArtifactIdentifier(
            conversation_id=conversation_id,
            patient_id=patient_id,
            filename=OUTPUT_PPTX_FILENAME.format(patient_id),
        )
        blob_path = self.data_access.chat_artifact_accessor.get_blob_path(artifact_id)
        output_url = get_chat_artifacts_url(blob_path)

        artifact = ChatArtifact(artifact_id=artifact_id, data=stream.getvalue())
        await self.data_access.chat_artifact_accessor.write(artifact)

        return (
            f"The PowerPoint presentation has been successfully created. "
            f'You can download it using the link below:<br><br>'
            f'<a href="{output_url}">{artifact_id.filename}</a>'
        )

    # ── Slide builders ──

    def _build_slide_1(self, slide, sc: SlideContent):
        """Patient Overview slide."""
        for shape in slide.shapes:
            if shape.name == "title":
                self._set_text(shape, sc.overview_title, size=32, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            elif shape.name == "subtitle":
                self._set_text(shape, sc.overview_subtitle, size=18, color=TEAL)
            elif shape.name == "body":
                self._set_bullets(shape, sc.overview_bullets)

    def _build_slide_2(self, slide, sc: SlideContent, chart_buf: Optional[BytesIO]):
        """Clinical Findings slide with tumor marker chart."""
        for shape in slide.shapes:
            if shape.name == "title":
                self._set_text(shape, sc.findings_title, size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            elif shape.name == "body_left":
                self._set_bullets(shape, sc.findings_bullets)
            elif shape.name == "chart_title":
                self._set_text(shape, sc.findings_chart_title, size=14, bold=True, color=TEAL)

        # Replace chart placeholder with actual chart image
        if chart_buf:
            # Find and remove the gray placeholder
            chart_placeholder = None
            for shape in slide.shapes:
                if shape.name == "chart_area":
                    chart_placeholder = shape
                    break

            if chart_placeholder:
                left = chart_placeholder.left
                top = chart_placeholder.top
                width = chart_placeholder.width
                height = chart_placeholder.height
                # Remove placeholder shape
                sp = chart_placeholder._element
                sp.getparent().remove(sp)
                # Add chart image
                slide.shapes.add_picture(chart_buf, left, top, width, height)

    def _build_slide_3(self, slide, sc: SlideContent):
        """Treatment & Trials slide."""
        for shape in slide.shapes:
            if shape.name == "title":
                self._set_text(shape, sc.treatment_title, size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            elif shape.name == "body":
                self._set_bullets(shape, sc.treatment_bullets)
            elif shape.name == "trials_body":
                self._set_bullets(shape, sc.trial_entries, color=NAVY)

    # ── Helpers ──

    @staticmethod
    def _set_text(shape, text: str, size: int = 14, bold: bool = False, color=DARK_TEXT):
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    @staticmethod
    def _set_bullets(shape, items: list[str], color=BULLET_COLOR):
        tf = shape.text_frame
        tf.clear()
        for i, item in enumerate(items[:6]):  # hard cap at 6
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            run.font.size = Pt(14)
            run.font.color.rgb = color

    def _create_marker_chart(self, tumor_markers_str: str) -> Optional[BytesIO]:
        """Generate a tumor marker trend chart as PNG in BytesIO.

        Expects tumor_markers_str to be JSON with a list of
        {date, value, marker_name} entries, or a text summary.
        Returns None if no plottable data.
        """
        if not tumor_markers_str:
            return None

        # Try parsing as JSON
        try:
            data = json.loads(tumor_markers_str)
        except (json.JSONDecodeError, TypeError):
            return None

        # Support both list format and dict-with-markers format
        if isinstance(data, dict):
            # Try extracting from common formats
            markers = data.get("markers", data.get("results", []))
            if not markers:
                return None
            data = markers

        if not isinstance(data, list) or len(data) == 0:
            return None

        # Extract dates and values
        dates = []
        values = []
        for entry in data:
            date_str = entry.get("date", entry.get("OrderDate", entry.get("ResultDate", "")))
            val_str = entry.get("value", entry.get("ResultValue", entry.get("result_value", "")))
            if date_str and val_str:
                try:
                    values.append(float(str(val_str).replace(",", "")))
                    dates.append(date_str)
                except (ValueError, TypeError):
                    continue

        if len(dates) < 2:
            return None

        # Get reference range if available
        ref_upper = None
        marker_name = data[0].get("marker_name", data[0].get("ComponentName", "Tumor Marker"))
        if "ca-125" in marker_name.lower() or "ca125" in marker_name.lower():
            ref_upper = 35.0
        elif "he4" in marker_name.lower():
            ref_upper = 140.0
        elif "hcg" in marker_name.lower():
            ref_upper = 5.0

        # Create chart (try/finally ensures figure is always closed)
        fig = None
        try:
            fig, ax = plt.subplots(figsize=(6, 4))

            ax.plot(range(len(dates)), values, marker="o", color="#007C91",
                    linewidth=2, markersize=6, markerfacecolor="#1B365D")
            ax.set_xticks(range(len(dates)))
            ax.set_xticklabels(dates, rotation=45, ha="right", fontsize=8)
            ax.set_ylabel(f"{marker_name}", fontsize=10, color="#333333")
            ax.set_title(f"{marker_name} Trend", fontsize=12, fontweight="bold", color="#1B365D")

            if ref_upper is not None:
                ax.axhline(y=ref_upper, color="#CC0000", linestyle="--", linewidth=1,
                            label=f"Upper Normal ({ref_upper})")
                ax.legend(fontsize=8)

            ax.grid(axis="y", alpha=0.3)
            ax.spines["top"].set_visible(False)
            ax.spines["right"].set_visible(False)

            plt.tight_layout()

            buf = BytesIO()
            plt.savefig(buf, format="png", dpi=150, bbox_inches="tight")
            buf.seek(0)
            return buf
        finally:
            if fig is not None:
                plt.close(fig)

    async def _summarize_for_slides(self, all_data: dict) -> SlideContent:
        """Use LLM to summarize all agent data into SlideContent structure."""
        chat_history = ChatHistory()
        chat_history.add_system_message(SLIDE_SUMMARIZATION_PROMPT)
        chat_history.add_user_message(
            "Agent outputs for slide generation:\n" + json.dumps(all_data, indent=2, default=str)
        )

        if model_supports_temperature():
            settings = AzureChatPromptExecutionSettings(
                temperature=0.0, response_format=SlideContent
            )
        else:
            settings = AzureChatPromptExecutionSettings(response_format=SlideContent)

        chat_service = self.kernel.get_service(service_id="default")
        response = await chat_service.get_chat_message_content(
            chat_history=chat_history, settings=settings
        )

        try:
            parsed = json.loads(response.content)
            return SlideContent(**parsed)
        except Exception:
            logger.warning("LLM response did not match SlideContent schema, using fallback")
            return SlideContent(
                overview_title=f"Patient {all_data.get('patient_id', 'Unknown')} — {all_data.get('cancer_type', '')}",
                overview_subtitle=f"FIGO {all_data.get('figo_stage', 'N/A')} | {all_data.get('molecular_profile', '')}",
                overview_bullets=[
                    f"Age: {all_data.get('patient_age', 'N/A')}, Gender: {all_data.get('patient_gender', 'N/A')}",
                    f"Cancer: {all_data.get('cancer_type', 'N/A')}",
                    f"FIGO Stage: {all_data.get('figo_stage', 'N/A')}",
                    f"Molecular: {all_data.get('molecular_profile', 'N/A')}",
                ],
                findings_title="Pathology & Imaging Findings",
                findings_bullets=[
                    all_data.get("pathology_findings", "No pathology data")[:80],
                    all_data.get("radiology_findings", "No radiology data")[:80],
                ],
                findings_chart_title="Tumor Marker Trend",
                treatment_title="Treatment Plan & Clinical Trials",
                treatment_bullets=[
                    all_data.get("treatment_plan", "No treatment plan")[:80],
                    all_data.get("board_discussion", "No discussion")[:80],
                ],
                trial_entries=[all_data.get("clinical_trials", "No trials identified")[:80]],
            )
