#!/usr/bin/env python3
"""Generate the GYN Tumor Board PPTX template programmatically.

Produces tumor_board_slides.pptx with 5 slides — one per tumor board column:
  Slide 1 — Patient          (Col 0: case #, MRN, attending, RTC, location, path)
  Slide 2 — Diagnosis        (Col 1: narrative + staging/genetics in red)
  Slide 3 — Previous Tx      (Col 2: treatment history + CA-125 chart)
  Slide 4 — Imaging          (Col 3: dated imaging studies)
  Slide 5 — Discussion       (Col 4: review types, trial eligibility, plan)

Run: python scripts/generate_pptx_template.py
Output: src/scenarios/default/templates/tumor_board_slides.pptx
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Medical color palette
NAVY = RGBColor(0x1B, 0x36, 0x5D)
TEAL = RGBColor(0x00, 0x7C, 0x91)
RED = RGBColor(0xFF, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
FOOTER_GRAY = RGBColor(0x99, 0x99, 0x99)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FOOTER_Y = Inches(7.0)
FOOTER_H = Inches(0.4)
MARGIN_L = Inches(0.8)
CONTENT_W = Inches(11.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, height=Inches(1.2)):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_WIDTH, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.name = "header_bar"
    return bar


def _add_textbox(slide, left, top, width, height, text="", font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, name=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    if name:
        txBox.name = name
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def _add_teal_divider(slide, top):
    div = slide.shapes.add_shape(1, MARGIN_L, top, Inches(1.5), Inches(0.05))
    div.fill.solid()
    div.fill.fore_color.rgb = TEAL
    div.line.fill.background()


def _add_footer(slide):
    _add_textbox(slide, MARGIN_L, FOOTER_Y, CONTENT_W, FOOTER_H,
                 text="GYN Oncology Tumor Board", font_size=10,
                 color=FOOTER_GRAY, alignment=PP_ALIGN.CENTER, name="footer")


def _add_slide_number_badge(slide, number, total=5):
    """Small teal badge showing slide number (e.g. '2 / 5') in top-right corner."""
    badge_w = Inches(1.2)
    badge = slide.shapes.add_textbox(
        SLIDE_WIDTH - badge_w - Inches(0.1), Inches(0.1), badge_w, Inches(0.35)
    )
    badge.name = "slide_badge"
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.RIGHT


def create_template():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ── Slide 1: Patient (Col 0) ──────────────────────────────────────────────
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s1, WHITE)
    _add_header_bar(s1, height=Inches(1.2))
    _add_slide_number_badge(s1, 1)

    _add_textbox(s1, MARGIN_L, Inches(0.25), Inches(10), Inches(0.6),
                 text="Patient", font_size=32, bold=True, color=WHITE, name="title")
    _add_textbox(s1, MARGIN_L, Inches(0.85), Inches(10), Inches(0.3),
                 text="Col 0 — Case logistics", font_size=13, color=TEAL, name="subtitle")

    _add_teal_divider(s1, Inches(1.35))

    _add_textbox(s1, MARGIN_L, Inches(1.55), CONTENT_W, Inches(5.2),
                 text="", font_size=16, name="body")
    _add_footer(s1)

    # ── Slide 2: Diagnosis & Pertinent History (Col 1) ────────────────────────
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s2, WHITE)
    _add_header_bar(s2, height=Inches(1.2))
    _add_slide_number_badge(s2, 2)

    _add_textbox(s2, MARGIN_L, Inches(0.25), Inches(10), Inches(0.6),
                 text="Diagnosis & Pertinent History", font_size=28, bold=True,
                 color=WHITE, name="title")
    _add_textbox(s2, MARGIN_L, Inches(0.85), Inches(10), Inches(0.3),
                 text="Col 1 — Narrative + staging", font_size=13, color=TEAL, name="subtitle")

    _add_teal_divider(s2, Inches(1.35))

    # Narrative bullets (left, ~60% width)
    _add_textbox(s2, MARGIN_L, Inches(1.55), Inches(7.5), Inches(3.8),
                 text="", font_size=14, name="body")

    # Staging block in red (right-bottom area)
    _add_textbox(s2, MARGIN_L, Inches(5.5), Inches(11.5), Inches(0.35),
                 text="Primary Site:", font_size=13, bold=False, color=RED, name="staging_label")
    _add_textbox(s2, MARGIN_L, Inches(5.5), CONTENT_W, Inches(1.35),
                 text="", font_size=13, color=RED, name="staging_body")

    _add_footer(s2)

    # ── Slide 3: Previous Tx or Operative Findings, Tumor Markers (Col 2) ─────
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s3, WHITE)
    _add_header_bar(s3, height=Inches(1.2))
    _add_slide_number_badge(s3, 3)

    _add_textbox(s3, MARGIN_L, Inches(0.25), Inches(10), Inches(0.6),
                 text="Previous Tx or Operative Findings", font_size=26, bold=True,
                 color=WHITE, name="title")
    _add_textbox(s3, MARGIN_L, Inches(0.85), Inches(10), Inches(0.3),
                 text="Col 2 — Treatment history & tumor markers", font_size=13,
                 color=TEAL, name="subtitle")

    _add_teal_divider(s3, Inches(1.35))

    # Left column — treatment history bullets
    _add_textbox(s3, MARGIN_L, Inches(1.55), Inches(6.0), Inches(5.2),
                 text="", font_size=14, name="body_left")

    # Right column — chart area
    _add_textbox(s3, Inches(7.5), Inches(1.55), Inches(5.3), Inches(0.4),
                 text="CA-125 Trend", font_size=14, bold=True,
                 color=TEAL, alignment=PP_ALIGN.CENTER, name="chart_title")

    chart_ph = s3.shapes.add_shape(1, Inches(7.5), Inches(2.05), Inches(5.3), Inches(4.5))
    chart_ph.fill.solid()
    chart_ph.fill.fore_color.rgb = LIGHT_GRAY
    chart_ph.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    chart_ph.name = "chart_area"

    _add_footer(s3)

    # ── Slide 4: Imaging (Col 3) ──────────────────────────────────────────────
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s4, WHITE)
    _add_header_bar(s4, height=Inches(1.2))
    _add_slide_number_badge(s4, 4)

    _add_textbox(s4, MARGIN_L, Inches(0.25), Inches(10), Inches(0.6),
                 text="Imaging", font_size=32, bold=True, color=WHITE, name="title")
    _add_textbox(s4, MARGIN_L, Inches(0.85), Inches(10), Inches(0.3),
                 text="Col 3 — Dated imaging studies", font_size=13, color=TEAL, name="subtitle")

    _add_teal_divider(s4, Inches(1.35))

    # Full-width body (imaging studies are verbose)
    _add_textbox(s4, MARGIN_L, Inches(1.55), CONTENT_W, Inches(5.2),
                 text="", font_size=13, name="body")

    _add_footer(s4)

    # ── Slide 5: Discussion (Col 4) ───────────────────────────────────────────
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s5, WHITE)
    _add_header_bar(s5, height=Inches(1.2))
    _add_slide_number_badge(s5, 5)

    _add_textbox(s5, MARGIN_L, Inches(0.25), Inches(10), Inches(0.6),
                 text="Discussion", font_size=32, bold=True, color=WHITE, name="title")
    _add_textbox(s5, MARGIN_L, Inches(0.85), Inches(10), Inches(0.3),
                 text="Col 4 — Review · Trial eligibility · Plan", font_size=13,
                 color=TEAL, name="subtitle")

    _add_teal_divider(s5, Inches(1.35))

    # Review types + trial eligibility (upper area)
    _add_textbox(s5, MARGIN_L, Inches(1.55), CONTENT_W, Inches(0.9),
                 text="", font_size=14, name="review_header")

    # Plan / action bullets (middle)
    _add_textbox(s5, MARGIN_L, Inches(2.55), CONTENT_W, Inches(2.8),
                 text="", font_size=14, name="body")

    # Trials section
    _add_textbox(s5, MARGIN_L, Inches(5.45), CONTENT_W, Inches(0.4),
                 text="Eligible Clinical Trials", font_size=16, bold=True,
                 color=TEAL, name="trials_header")
    _add_textbox(s5, MARGIN_L, Inches(5.95), CONTENT_W, Inches(0.9),
                 text="", font_size=13, name="trials_body")

    _add_footer(s5)

    return prs


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(script_dir)
    output_path = os.path.join(
        repo_root, "src", "scenarios", "default", "templates", "tumor_board_slides.pptx"
    )

    prs = create_template()
    prs.save(output_path)
    print(f"Template saved to: {output_path}")


if __name__ == "__main__":
    main()
